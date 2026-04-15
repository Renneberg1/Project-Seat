"""Add a Project Spin-Up slide to the existing presentation.

Inserts after slide 3 (Meet Project Seat divider), before the workflow slide.
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os, copy

# ── Brand colors ──────────────────────────────────────────────────────────
NAVY = RGBColor(0x00, 0x0F, 0x5B)
ROYAL = RGBColor(0x09, 0x72, 0xF1)
SKY = RGBColor(0x70, 0xD4, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF1, 0xF1, 0xF1)
DARK_GRAY = RGBColor(0x89, 0x89, 0x89)
MED_GRAY = RGBColor(0x9D, 0x9D, 0x9C)
LIGHT_GRAY = RGBColor(0xD0, 0xD0, 0xD0)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREEN = RGBColor(0x34, 0xA8, 0x53)
AMBER = RGBColor(0xF9, 0xAB, 0x00)
PURPLE = RGBColor(0x93, 0x34, 0xE6)

FONT = "Rules"
FONT_MED = "Rules Medium"

SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

PPTX_PATH = os.path.join(os.path.dirname(__file__), "..", "Project Seat Presentation.pptx")


def add_rich_textbox(slide, left, top, width, height, lines, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    for i, ld in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ld.get("alignment", PP_ALIGN.LEFT)
        if "space_before" in ld:
            p.space_before = ld["space_before"]
        if "space_after" in ld:
            p.space_after = ld["space_after"]
        run = p.add_run()
        run.text = ld["text"]
        run.font.name = ld.get("font_name", FONT)
        run.font.size = ld.get("font_size", Pt(14))
        run.font.color.rgb = ld.get("color", BLACK)
        if ld.get("bold"):
            run.font.bold = True
    return txBox


def add_textbox(slide, left, top, width, height, text,
                font_name=FONT, font_size=Pt(14), color=BLACK,
                bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox


def add_rounded_box(slide, left, top, width, height, fill_color, border_color=None,
                    text="", font_name=FONT, font_size=Pt(12), text_color=WHITE,
                    bold=False, alignment=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.1
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = text_color
    run.font.bold = bold
    from pptx.oxml.ns import qn
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')
    return shape


def add_chevron_arrow(slide, left, top, width, height, color=ROYAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_down_arrow_shape(slide, left, top, width, height, color=ROYAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def clear_placeholder(ph):
    for p in ph.text_frame.paragraphs:
        p.clear()


def set_text(ph, text, font_name=None, font_size=None, color=None, bold=None):
    tf = ph.text_frame
    for p in tf.paragraphs:
        p.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    if font_name:
        run.font.name = font_name
    if font_size:
        run.font.size = font_size
    if color:
        run.font.color.rgb = color
    if bold is not None:
        run.font.bold = bold


def move_slide(prs, from_index, to_index):
    """Move a slide from from_index to to_index (0-based)."""
    sldIdLst = prs.slides._sldIdLst
    el = sldIdLst[from_index]
    sldIdLst.remove(el)
    if to_index >= len(sldIdLst):
        sldIdLst.append(el)
    else:
        sldIdLst.insert(to_index, el)


def main():
    prs = Presentation(PPTX_PATH)

    print(f"Current slides: {len(prs.slides)}")

    # Use Content (1 Column) layout — index 8
    layout = prs.slide_layouts[8]
    slide = prs.slides.add_slide(layout)

    # Set section label and title via placeholders
    set_text(slide.placeholders[13], "PROJECT SPIN-UP",
             font_name=FONT_MED, font_size=Pt(10), color=ROYAL)
    set_text(slide.placeholders[0], "One-Click Project Creation",
             font_name=FONT_MED, font_size=Pt(36), color=NAVY)
    # Clear body placeholders — we use shapes
    for idx in [2, 14]:
        if idx in slide.placeholders:
            clear_placeholder(slide.placeholders[idx])

    # ── Visual: Form → Approval Queue → Created Assets ──

    # Row 1: Form input (left) → arrow → Approval queue (right)
    row1_y = Emu(1800000)
    box_h = Emu(1500000)

    # Left: Form
    form_x = Emu(300000)
    form_w = Emu(3600000)
    add_rounded_box(slide, form_x, row1_y, form_w, box_h, OFF_WHITE,
                    border_color=ROYAL, text="", text_color=BLACK)
    add_rich_textbox(slide, form_x + Emu(120000), row1_y + Emu(80000),
                     form_w - Emu(240000), box_h - Emu(160000), [
        {"text": "SPIN-UP FORM", "font_name": FONT_MED, "font_size": Pt(12),
         "color": ROYAL, "bold": True, "space_after": Pt(8)},
        {"text": "\u2022  Project name & program", "font_name": FONT, "font_size": Pt(10), "color": BLACK, "space_before": Pt(3)},
        {"text": "\u2022  Team projects (AIM, CTCV...)", "font_name": FONT, "font_size": Pt(10), "color": BLACK, "space_before": Pt(3)},
        {"text": "\u2022  Target date & labels", "font_name": FONT, "font_size": Pt(10), "color": BLACK, "space_before": Pt(3)},
        {"text": "\u2022  Goal description", "font_name": FONT, "font_size": Pt(10), "color": BLACK, "space_before": Pt(3)},
    ])

    # Arrow
    arr1_x = form_x + form_w + Emu(60000)
    add_chevron_arrow(slide, arr1_x, row1_y + box_h // 2 - Emu(150000),
                      Emu(300000), Emu(300000), color=ROYAL)

    # Middle: Approval Queue
    queue_x = arr1_x + Emu(420000)
    queue_w = Emu(3600000)
    add_rounded_box(slide, queue_x, row1_y, queue_w, box_h, AMBER,
                    text="", text_color=NAVY)
    add_rich_textbox(slide, queue_x + Emu(120000), row1_y + Emu(80000),
                     queue_w - Emu(240000), box_h - Emu(160000), [
        {"text": "APPROVAL QUEUE", "font_name": FONT_MED, "font_size": Pt(12),
         "color": NAVY, "bold": True, "space_after": Pt(8)},
        {"text": "\u2022  4\u20136+ items queued", "font_name": FONT, "font_size": Pt(10), "color": NAVY, "space_before": Pt(3)},
        {"text": "\u2022  Review each action", "font_name": FONT, "font_size": Pt(10), "color": NAVY, "space_before": Pt(3)},
        {"text": "\u2022  Approve one-by-one", "font_name": FONT, "font_size": Pt(10), "color": NAVY, "space_before": Pt(3)},
        {"text": "     or approve all at once", "font_name": FONT, "font_size": Pt(10), "color": NAVY, "space_before": Pt(1)},
    ])

    # Arrow
    arr2_x = queue_x + queue_w + Emu(60000)
    add_chevron_arrow(slide, arr2_x, row1_y + box_h // 2 - Emu(150000),
                      Emu(300000), Emu(300000), color=GREEN)

    # Right: Executed
    exec_x = arr2_x + Emu(420000)
    exec_w = Emu(3600000)
    add_rounded_box(slide, exec_x, row1_y, exec_w, box_h, GREEN,
                    text="", text_color=WHITE)
    add_rich_textbox(slide, exec_x + Emu(120000), row1_y + Emu(80000),
                     exec_w - Emu(240000), box_h - Emu(160000), [
        {"text": "EXECUTED", "font_name": FONT_MED, "font_size": Pt(12),
         "color": WHITE, "bold": True, "space_after": Pt(8)},
        {"text": "\u2022  Sequential execution", "font_name": FONT, "font_size": Pt(10), "color": WHITE, "space_before": Pt(3)},
        {"text": "\u2022  Sentinel dependencies", "font_name": FONT, "font_size": Pt(10), "color": WHITE, "space_before": Pt(3)},
        {"text": "     auto-resolved", "font_name": FONT, "font_size": Pt(10), "color": WHITE, "space_before": Pt(1)},
        {"text": "\u2022  Full audit trail", "font_name": FONT, "font_size": Pt(10), "color": WHITE, "space_before": Pt(3)},
    ])

    # ── Row 2: Created assets (what gets built) ──
    row2_y = Emu(3700000)
    asset_h = Emu(2600000)
    asset_gap = Emu(200000)

    # 5 asset cards showing what gets created
    assets = [
        ("\u2460", "Jira Goal", "PROG project\nGoal ticket with\nADF description\n+ labels + due date", NAVY),
        ("\u2461", "Risk Version", "RISK project\nFix version for\nrisk tracking\nlinked to project", ROYAL),
        ("\u2462", "Team Versions", "Per-team fix\nversions (AIM,\nCTCV, YAM...)\none per team", ROYAL),
        ("\u2463", "Charter Page", "Confluence page\nfrom template with\nplaceholders\nreplaced", PURPLE),
        ("\u2464", "XFT Page", "Child of Charter\npage from template\n+ Goal description\nupdated with links", PURPLE),
    ]

    asset_total_w = SLIDE_W - Emu(600000)
    asset_w = (asset_total_w - 4 * asset_gap) // 5
    asset_start_x = Emu(300000)

    for i, (num, title, desc, accent) in enumerate(assets):
        x = asset_start_x + i * (asset_w + asset_gap)

        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, row2_y, asset_w, asset_h)
        card.fill.solid()
        card.fill.fore_color.rgb = OFF_WHITE
        card.line.fill.background()
        card.adjustments[0] = 0.06

        # Accent bar top
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, row2_y, asset_w, Emu(60000))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        # Number circle + title
        add_rich_textbox(slide, x + Emu(80000), row2_y + Emu(120000),
                         asset_w - Emu(160000), asset_h - Emu(200000), [
            {"text": num, "font_name": FONT_MED, "font_size": Pt(20),
             "color": accent, "bold": True, "space_after": Pt(2)},
            {"text": title, "font_name": FONT_MED, "font_size": Pt(13),
             "color": accent, "bold": True, "space_before": Pt(2)},
            {"text": desc, "font_name": FONT, "font_size": Pt(10),
             "color": DARK_GRAY, "space_before": Pt(8)},
        ])

    # ── Caption at bottom ──
    add_rich_textbox(slide, Emu(300000), Emu(6400000), Emu(11500000), Emu(350000), [
        {"text": "One form submission \u2192 6+ coordinated actions across Jira & Confluence \u2192 project ready in seconds",
         "font_name": FONT, "font_size": Pt(12), "color": DARK_GRAY,
         "alignment": PP_ALIGN.CENTER},
    ])

    # ── Move the new slide (currently last) to position 4 (0-indexed = 3) ──
    # After "Meet Project Seat" (slide 3) and before "End-to-End Workflow" (slide 4)
    current_idx = len(prs.slides) - 1
    target_idx = 3  # insert as slide 4 (0-based index 3)
    move_slide(prs, current_idx, target_idx)

    prs.save(PPTX_PATH)
    print(f"Saved. Total slides: {len(prs.slides)}")
    print(f"New slide inserted at position {target_idx + 1}")

    # Verify order
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = p.text.strip()
                    if t and len(t) > 5:
                        texts.append(t[:50])
                        break
        print(f"  Slide {i+1}: {texts[:2]}")


if __name__ == "__main__":
    main()
