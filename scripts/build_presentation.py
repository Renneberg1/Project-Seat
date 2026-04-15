"""Build the Project Seat presentation with visual-first slides."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand colors ──────────────────────────────────────────────────────────
NAVY = RGBColor(0x00, 0x0F, 0x5B)
ROYAL = RGBColor(0x09, 0x72, 0xF1)
SKY = RGBColor(0x70, 0xD4, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF1, 0xF1, 0xF1)
DARK_GRAY = RGBColor(0x89, 0x89, 0x89)
MED_GRAY = RGBColor(0x9D, 0x9D, 0x9C)
LIGHT_GRAY = RGBColor(0xD0, 0xD0, 0xD0)
BLACK = RGBColor(0x00, 0x00, 0x00)
# Accent colors for diagrams
GREEN = RGBColor(0x34, 0xA8, 0x53)
AMBER = RGBColor(0xF9, 0xAB, 0x00)
RED_ACCENT = RGBColor(0xEA, 0x43, 0x35)
PURPLE = RGBColor(0x93, 0x34, 0xE6)

FONT = "Rules"
FONT_MED = "Rules Medium"

PPTX_PATH = os.path.join(os.path.dirname(__file__), "..", "Project Seat Presentation.pptx")

# Slide dimensions (from template)
SLIDE_W = Emu(12192000)  # 13.33 inches
SLIDE_H = Emu(6858000)   # 7.50 inches


# ── Helper functions ──────────────────────────────────────────────────────

def clear_placeholder(ph):
    for p in ph.text_frame.paragraphs:
        p.clear()


def set_text(ph, text, font_name=None, font_size=None, color=None, bold=None, alignment=None):
    tf = ph.text_frame
    for p in tf.paragraphs:
        p.clear()
    p = tf.paragraphs[0]
    if alignment:
        p.alignment = alignment
    run = p.add_run()
    run.text = text
    if font_name:
        run.font.name = font_name
    if font_size:
        run.font.size = font_size
    if color:
        run.font.color.rgb = color
    if bold is not None:
        run.font.bold = bold


def add_textbox(slide, left, top, width, height, text,
                font_name=FONT, font_size=Pt(14), color=BLACK,
                bold=False, alignment=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox


def add_rich_textbox(slide, left, top, width, height, lines, word_wrap=True):
    """lines = list of dicts: text, font_name, font_size, color, bold, alignment, space_before"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    for i, ld in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ld.get("alignment", PP_ALIGN.LEFT)
        if "space_before" in ld:
            p.space_before = ld["space_before"]
        if "space_after" in ld:
            p.space_after = ld["space_after"]
        run = p.add_run()
        run.text = ld["text"]
        run.font.name = ld.get("font_name", FONT)
        run.font.size = ld.get("font_size", Pt(14))
        run.font.color.rgb = ld.get("color", BLACK)
        if ld.get("bold"):
            run.font.bold = True
    return txBox


def add_rounded_box(slide, left, top, width, height, fill_color, border_color=None,
                    text="", font_name=FONT, font_size=Pt(12), text_color=WHITE,
                    bold=False, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    """Add a rounded rectangle with text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    # Adjust corner radius
    shape.adjustments[0] = 0.1
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = text_color
    run.font.bold = bold
    # Vertical centering
    from pptx.oxml.ns import qn
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')
    return shape


def add_arrow(slide, start_left, start_top, end_left, end_top, color=ROYAL, width=Pt(2)):
    """Add a connector arrow between two points."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_left, start_top, end_left, end_top
    )
    connector.line.color.rgb = color
    connector.line.width = width
    # Add arrowhead
    from pptx.oxml.ns import qn
    ln = connector.line._ln
    tailEnd = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tailEnd)
    return connector


def add_chevron_arrow(slide, left, top, width, height, color=ROYAL):
    """Add a right-pointing chevron/arrow shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_screenshot_placeholder(slide, left, top, width, height, label="Screenshot"):
    """Add a dashed-border placeholder rectangle for a screenshot."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = OFF_WHITE
    shape.line.color.rgb = LIGHT_GRAY
    shape.line.width = Pt(2)
    shape.line.dash_style = 2  # dash
    shape.adjustments[0] = 0.02
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"\n\n{label}\n(drop screenshot here)"
    run.font.name = FONT
    run.font.size = Pt(11)
    run.font.color.rgb = MED_GRAY
    from pptx.oxml.ns import qn
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')
    return shape


def add_bullet_text(ph, items, font_size=Pt(14), color=BLACK, spacing=Pt(6), bold_first=False):
    tf = ph.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = spacing
        if bold_first and " — " in item:
            bp, rest = item.split(" — ", 1)
            r1 = p.add_run()
            r1.text = bp + " — "
            r1.font.name = FONT_MED
            r1.font.size = font_size
            r1.font.color.rgb = color
            r1.font.bold = True
            r2 = p.add_run()
            r2.text = rest
            r2.font.name = FONT
            r2.font.size = font_size
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.name = FONT
            r.font.size = font_size
            r.font.color.rgb = color


# ── Build the deck ────────────────────────────────────────────────────────

def build_presentation():
    prs = Presentation(PPTX_PATH)

    # Remove all existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId is None:
            rId = prs.slides._sldIdLst[0].attrib.get('r:id')
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # Layout references
    cover1 = prs.slide_layouts[0]       # Cover 1
    divider1 = prs.slide_layouts[2]     # Divider 1 (white bg)
    divider2 = prs.slide_layouts[3]     # Divider 2 (navy bg)
    content1 = prs.slide_layouts[8]     # Content (1 Column)
    content2 = prs.slide_layouts[9]     # Content (2 Column)
    content3 = prs.slide_layouts[10]    # Content (3 Column)
    color_blue = prs.slide_layouts[12]  # Content Color Blue
    color_navy = prs.slide_layouts[13]  # Content Color Navy
    bg_navy = prs.slide_layouts[14]     # Color Background - Navy
    bg_royal = prs.slide_layouts[15]    # Color Background - Royal

    # =====================================================================
    # SLIDE 1 — Cover
    # =====================================================================
    slide = prs.slides.add_slide(cover1)
    for ph in slide.placeholders:
        clear_placeholder(ph)
    add_textbox(slide, Emu(505141), Emu(2600000), Emu(8500000), Emu(1800000),
                "PROJECT SEAT",
                font_name=FONT_MED, font_size=Pt(60), color=NAVY, bold=True)
    add_textbox(slide, Emu(505141), Emu(4400000), Emu(8500000), Emu(800000),
                "AI-Powered Project Management\nfor Medical Device Engineering",
                font_name=FONT, font_size=Pt(22), color=DARK_GRAY)
    set_text(slide.placeholders[13], "MARCH 2026  |  THOMAS RENNEBERG",
             font_name=FONT, font_size=Pt(10), color=MED_GRAY)
    set_text(slide.placeholders[12], "AI SOLUTIONS CHALLENGE",
             font_name=FONT, font_size=Pt(10), color=MED_GRAY)

    # =====================================================================
    # SLIDE 2 — The Problem (divider)
    # =====================================================================
    slide = prs.slides.add_slide(divider2)
    set_text(slide.placeholders[0], "Medical Device PM\nis Drowning\nin Process",
             font_name=FONT_MED, font_size=Pt(50), color=WHITE)
    set_text(slide.placeholders[13],
             "Regulated environments demand rigorous traceability — but existing tools add overhead, not intelligence.",
             font_name=FONT, font_size=Pt(12), color=SKY)

    # =====================================================================
    # SLIDE 3 — The Problem — Visual with pain points
    # =====================================================================
    slide = prs.slides.add_slide(content1)
    set_text(slide.placeholders[13], "THE CHALLENGE",
             font_name=FONT_MED, font_size=Pt(10), color=ROYAL)
    set_text(slide.placeholders[0], "Where Time Goes Today",
             font_name=FONT_MED, font_size=Pt(36), color=NAVY)
    # Hide the content placeholders — we'll use shapes
    for idx in [2, 14]:
        if idx in slide.placeholders:
            clear_placeholder(slide.placeholders[idx])

    # Four pain-point boxes arranged as 2x2 grid
    box_w = Emu(4800000)
    box_h = Emu(1600000)
    gap_x = Emu(300000)
    gap_y = Emu(250000)
    start_x = Emu(400000)
    start_y = Emu(2200000)

    pain_points = [
        ("Meeting \u2192 Jira Gap", "Risks and decisions discussed in meetings\nnever make it into Jira. Institutional\nknowledge evaporates after every call.", RED_ACCENT),
        ("Stale Documentation", "Charter and compliance documents\ngo out of date within weeks.\nNo one knows the current state.", AMBER),
        ("Manual Reporting", "CEO status updates and health reviews\ntake days of manual data gathering\nacross Jira, Confluence, and spreadsheets.", PURPLE),
        ("Invisible Project Health", "Problems are invisible until they're\ncritical. No early warning system.\nNo structured assessment.", NAVY),
    ]
    for i, (title, desc, accent) in enumerate(pain_points):
        col = i % 2
        row = i // 2
        x = start_x + col * (box_w + gap_x)
        y = start_y + row * (box_h + gap_y)

        # Accent bar on left
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Emu(60000), box_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        # Text box
        add_rich_textbox(slide, x + Emu(120000), y + Emu(100000),
                         box_w - Emu(180000), box_h - Emu(200000), [
            {"text": title, "font_name": FONT_MED, "font_size": Pt(16),
             "color": accent, "bold": True, "space_after": Pt(6)},
            {"text": desc, "font_name": FONT, "font_size": Pt(12),
             "color": DARK_GRAY, "space_before": Pt(4)},
        ])

    # =====================================================================
    # SLIDE 4 — The Solution (divider)
    # =====================================================================
    slide = prs.slides.add_slide(divider2)
    set_text(slide.placeholders[0], "Meet\nProject Seat",
             font_name=FONT_MED, font_size=Pt(60), color=WHITE)
    set_text(slide.placeholders[13],
             "An AI cockpit that connects Jira, Confluence, and Zoom — with an LLM as the reasoning engine.",
             font_name=FONT, font_size=Pt(12), color=SKY)

    # =====================================================================
    # SLIDE 5 — Product Workflow (visual flow diagram)
    # =====================================================================
    slide = prs.slides.add_slide(content1)
    set_text(slide.placeholders[13], "HOW IT WORKS",
             font_name=FONT_MED, font_size=Pt(10), color=ROYAL)
    set_text(slide.placeholders[0], "End-to-End Workflow",
             font_name=FONT_MED, font_size=Pt(36), color=NAVY)
    for idx in [2, 14]:
        if idx in slide.placeholders:
            clear_placeholder(slide.placeholders[idx])

    # Flow: Input → LLM Agents → Approval Gate → Outputs
    flow_y = Emu(2300000)
    box_h_flow = Emu(3600000)
    arrow_w = Emu(250000)
    margin = Emu(250000)

    # Calculate box widths to fill the slide
    usable = SLIDE_W - 2 * margin - 3 * arrow_w
    bw = usable // 4

    positions = []
    x = margin
    for i in range(4):
        positions.append(x)
        x += bw
        if i < 3:
            x += arrow_w

    # Box 1: Inputs
    add_rounded_box(slide, positions[0], flow_y, bw, box_h_flow, OFF_WHITE,
                    border_color=LIGHT_GRAY, text="", text_color=BLACK)
    add_rich_textbox(slide, positions[0] + Emu(80000), flow_y + Emu(100000),
                     bw - Emu(160000), box_h_flow - Emu(200000), [
        {"text": "DATA INPUTS", "font_name": FONT_MED, "font_size": Pt(11),
         "color": ROYAL, "bold": True, "space_after": Pt(12)},
        {"text": "\u2022  Meeting transcripts\n   (.vtt / .txt / .docx)",
         "font_name": FONT, "font_size": Pt(10), "color": BLACK, "space_before": Pt(6)},
        {"text": "\u2022  Zoom recordings\n   (auto-ingested via OAuth)",
         "font_name": FONT, "font_size": Pt(10), "color": BLACK, "space_before": Pt(6)},
        {"text": "\u2022  PM notes & context",
         "font_name": FONT, "font_size": Pt(10), "color": BLACK, "space_before": Pt(6)},
        {"text": "\u2022  Jira + Confluence\n   live project data",
         "font_name": FONT, "font_size": Pt(10), "color": BLACK, "space_before": Pt(6)},
    ])

    # Box 2: LLM Agents
    add_rounded_box(slide, positions[1], flow_y, bw, box_h_flow, ROYAL,
                    text="", text_color=WHITE)
    add_rich_textbox(slide, positions[1] + Emu(80000), flow_y + Emu(100000),
                     bw - Emu(160000), box_h_flow - Emu(200000), [
        {"text": "7 AI AGENTS", "font_name": FONT_MED, "font_size": Pt(11),
         "color": SKY, "bold": True, "space_after": Pt(12)},
        {"text": "\u2022  Transcript Analysis",
         "font_name": FONT, "font_size": Pt(10), "color": WHITE, "space_before": Pt(4)},
        {"text": "\u2022  Charter Update",
         "font_name": FONT, "font_size": Pt(10), "color": WHITE, "space_before": Pt(4)},
        {"text": "\u2022  Health Review",
         "font_name": FONT, "font_size": Pt(10), "color": WHITE, "space_before": Pt(4)},
        {"text": "\u2022  CEO Review",
         "font_name": FONT, "font_size": Pt(10), "color": WHITE, "space_before": Pt(4)},
        {"text": "\u2022  Closure Report",
         "font_name": FONT, "font_size": Pt(10), "color": WHITE, "space_before": Pt(4)},
        {"text": "\u2022  Risk Refinement",
         "font_name": FONT, "font_size": Pt(10), "color": WHITE, "space_before": Pt(4)},
        {"text": "\u2022  Zoom Matching",
         "font_name": FONT, "font_size": Pt(10), "color": WHITE, "space_before": Pt(4)},
    ])

    # Box 3: Approval Gate
    add_rounded_box(slide, positions[2], flow_y, bw, box_h_flow, AMBER,
                    text="", text_color=BLACK)
    add_rich_textbox(slide, positions[2] + Emu(80000), flow_y + Emu(100000),
                     bw - Emu(160000), box_h_flow - Emu(200000), [
        {"text": "APPROVAL GATE", "font_name": FONT_MED, "font_size": Pt(11),
         "color": NAVY, "bold": True, "space_after": Pt(12)},
        {"text": "\u2022  Human reviews every\n   AI suggestion",
         "font_name": FONT, "font_size": Pt(10), "color": NAVY, "space_before": Pt(6)},
        {"text": "\u2022  Accept / Reject / Edit\n   per item",
         "font_name": FONT, "font_size": Pt(10), "color": NAVY, "space_before": Pt(6)},
        {"text": "\u2022  Second approval\n   before execution",
         "font_name": FONT, "font_size": Pt(10), "color": NAVY, "space_before": Pt(6)},
        {"text": "\u2022  Full audit trail\n   in SQLite",
         "font_name": FONT, "font_size": Pt(10), "color": NAVY, "space_before": Pt(6)},
    ])

    # Box 4: Outputs
    add_rounded_box(slide, positions[3], flow_y, bw, box_h_flow, NAVY,
                    text="", text_color=WHITE)
    add_rich_textbox(slide, positions[3] + Emu(80000), flow_y + Emu(100000),
                     bw - Emu(160000), box_h_flow - Emu(200000), [
        {"text": "OUTPUTS", "font_name": FONT_MED, "font_size": Pt(11),
         "color": SKY, "bold": True, "space_after": Pt(12)},
        {"text": "\u2022  Jira risk tickets",
         "font_name": FONT, "font_size": Pt(10), "color": WHITE, "space_before": Pt(4)},
        {"text": "\u2022  Jira decision tickets",
         "font_name": FONT, "font_size": Pt(10), "color": WHITE, "space_before": Pt(4)},
        {"text": "\u2022  Charter page edits",
         "font_name": FONT, "font_size": Pt(10), "color": WHITE, "space_before": Pt(4)},
        {"text": "\u2022  XFT meeting minutes",
         "font_name": FONT, "font_size": Pt(10), "color": WHITE, "space_before": Pt(4)},
        {"text": "\u2022  CEO review pages",
         "font_name": FONT, "font_size": Pt(10), "color": WHITE, "space_before": Pt(4)},
        {"text": "\u2022  Closure report pages",
         "font_name": FONT, "font_size": Pt(10), "color": WHITE, "space_before": Pt(4)},
        {"text": "\u2022  Knowledge base entries",
         "font_name": FONT, "font_size": Pt(10), "color": WHITE, "space_before": Pt(4)},
    ])

    # Arrows between boxes
    arrow_y_mid = flow_y + box_h_flow // 2
    for i in range(3):
        arr_x_start = positions[i] + bw
        arr_x_end = positions[i + 1]
        add_chevron_arrow(slide, arr_x_start + Emu(20000), arrow_y_mid - Emu(180000),
                          arrow_w - Emu(40000), Emu(360000),
                          color=ROYAL if i < 2 else GREEN)

    # =====================================================================
    # SLIDE 6 — AI Meeting Intelligence (screenshot + explanation)
    # =====================================================================
    slide = prs.slides.add_slide(content1)
    set_text(slide.placeholders[13], "INNOVATION 1",
             font_name=FONT_MED, font_size=Pt(10), color=ROYAL)
    set_text(slide.placeholders[0], "AI Meeting Intelligence",
             font_name=FONT_MED, font_size=Pt(36), color=NAVY)
    for idx in [2, 14]:
        if idx in slide.placeholders:
            clear_placeholder(slide.placeholders[idx])

    # Left side: screenshot placeholder
    add_screenshot_placeholder(slide, Emu(300000), Emu(2100000),
                               Emu(5800000), Emu(4400000),
                               "Transcript Suggestions UI")

    # Right side: key points
    add_rich_textbox(slide, Emu(6400000), Emu(2100000), Emu(5400000), Emu(4400000), [
        {"text": "Upload or auto-ingest", "font_name": FONT_MED, "font_size": Pt(16),
         "color": NAVY, "bold": True},
        {"text": ".vtt, .txt, .docx transcripts — or Zoom\nrecordings fetched automatically via OAuth",
         "font_name": FONT, "font_size": Pt(12), "color": DARK_GRAY, "space_before": Pt(4)},
        {"text": "LLM extracts structured data", "font_name": FONT_MED, "font_size": Pt(16),
         "color": NAVY, "bold": True, "space_before": Pt(24)},
        {"text": "Risks, decisions, action items, document\nupdates — all mapped to Jira/Confluence",
         "font_name": FONT, "font_size": Pt(12), "color": DARK_GRAY, "space_before": Pt(4)},
        {"text": "ISO 14971 quality refinement", "font_name": FONT_MED, "font_size": Pt(16),
         "color": NAVY, "bold": True, "space_before": Pt(24)},
        {"text": "Multi-round Q&A loop refines risks against\nregulatory quality criteria (up to 5 rounds)",
         "font_name": FONT, "font_size": Pt(12), "color": DARK_GRAY, "space_before": Pt(4)},
        {"text": "Semantic dedup", "font_name": FONT_MED, "font_size": Pt(16),
         "color": NAVY, "bold": True, "space_before": Pt(24)},
        {"text": "Links findings to existing Jira items.\nNo duplicates. Context-enriched two-pass analysis.",
         "font_name": FONT, "font_size": Pt(12), "color": DARK_GRAY, "space_before": Pt(4)},
    ])

    # =====================================================================
    # SLIDE 7 — Approval Queue (screenshot + flow)
    # =====================================================================
    slide = prs.slides.add_slide(color_blue)
    set_text(slide.placeholders[13], "INNOVATION 2",
             font_name=FONT_MED, font_size=Pt(10), color=WHITE)
    set_text(slide.placeholders[0], "Human-in-the-Loop Governance",
             font_name=FONT_MED, font_size=Pt(36), color=WHITE)
    for idx in [2, 14]:
        if idx in slide.placeholders:
            clear_placeholder(slide.placeholders[idx])

    # Visual: 3-step approval flow
    step_w = Emu(3200000)
    step_h = Emu(1300000)
    step_y = Emu(2400000)
    step_gap = Emu(400000)
    total = 3 * step_w + 2 * step_gap
    sx = (SLIDE_W - total) // 2

    steps = [
        ("1. AI Suggests", "LLM generates structured\nrisks, decisions, edits\nfrom your data", WHITE, NAVY),
        ("2. Human Reviews", "Accept, reject, or refine\neach suggestion before\nit becomes an action", WHITE, ROYAL),
        ("3. Execute + Audit", "Approved actions written\nto Jira/Confluence with\nfull audit trail", WHITE, GREEN),
    ]
    for i, (title, desc, text_c, fill_c) in enumerate(steps):
        x = sx + i * (step_w + step_gap)
        add_rounded_box(slide, x, step_y, step_w, step_h, fill_c,
                        text="", text_color=text_c)
        add_rich_textbox(slide, x + Emu(120000), step_y + Emu(100000),
                         step_w - Emu(240000), step_h - Emu(200000), [
            {"text": title, "font_name": FONT_MED, "font_size": Pt(18),
             "color": text_c, "bold": True, "space_after": Pt(8)},
            {"text": desc, "font_name": FONT, "font_size": Pt(12),
             "color": SKY if fill_c == NAVY else WHITE},
        ])
        if i < 2:
            add_chevron_arrow(slide, x + step_w + Emu(40000),
                              step_y + step_h // 2 - Emu(150000),
                              step_gap - Emu(80000), Emu(300000), color=WHITE)

    # Bottom: screenshot placeholder
    add_screenshot_placeholder(slide, Emu(1500000), Emu(4100000),
                               Emu(9200000), Emu(2300000),
                               "Approval Queue UI")

    # =====================================================================
    # SLIDE 8 — Report Generation (3 report types visual)
    # =====================================================================
    slide = prs.slides.add_slide(content1)
    set_text(slide.placeholders[13], "INNOVATION 3",
             font_name=FONT_MED, font_size=Pt(10), color=ROYAL)
    set_text(slide.placeholders[0], "One-Click Structured Reports",
             font_name=FONT_MED, font_size=Pt(36), color=NAVY)
    for idx in [2, 14]:
        if idx in slide.placeholders:
            clear_placeholder(slide.placeholders[idx])

    # Shared pattern description
    add_textbox(slide, Emu(300000), Emu(1500000), Emu(11500000), Emu(500000),
                "Two-step Q&A \u2192 LLM asks clarifying questions, then generates hybrid data tables + narrative",
                font_name=FONT, font_size=Pt(13), color=DARK_GRAY)

    # Three report cards side by side
    card_w = Emu(3500000)
    card_h = Emu(4200000)
    card_y = Emu(2200000)
    card_gap = Emu(350000)
    total_cards = 3 * card_w + 2 * card_gap
    cx = (SLIDE_W - total_cards) // 2

    reports = [
        ("Health Review", GREEN, [
            "On-demand assessment",
            "Green / Amber / Red rating",
            "Pulls all project data:",
            "  risks, decisions, DHF,",
            "  team progress, charter",
            "Concerns + next actions",
            "",
            "Read-only — no approval",
            "needed",
        ]),
        ("CEO Review", ROYAL, [
            "Fortnightly executive update",
            "Last-2-weeks data lens",
            "Deterministic data tables:",
            "  milestones, risks, scope",
            "LLM-written commentary",
            "Health indicator + escalations",
            "",
            "Published to Confluence",
            "via approval queue",
        ]),
        ("Closure Report", NAVY, [
            "Full lifecycle summary",
            "Delivery outcome analysis",
            "Timeline + scope adherence",
            "Risk/issue closure status",
            "Success criteria assessment",
            "Lessons learned narrative",
            "",
            "New Confluence page",
            "via approval queue",
        ]),
    ]

    for i, (title, accent, bullets) in enumerate(reports):
        x = cx + i * (card_w + card_gap)

        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, card_y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = OFF_WHITE
        card.line.fill.background()
        card.adjustments[0] = 0.04

        # Accent bar top
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, card_y, card_w, Emu(80000))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        # Title
        add_textbox(slide, x + Emu(120000), card_y + Emu(150000),
                    card_w - Emu(240000), Emu(400000),
                    title, font_name=FONT_MED, font_size=Pt(18),
                    color=accent, bold=True)

        # Bullets
        lines = []
        for b in bullets:
            lines.append({"text": b, "font_name": FONT, "font_size": Pt(11),
                          "color": DARK_GRAY, "space_before": Pt(3)})
        add_rich_textbox(slide, x + Emu(120000), card_y + Emu(600000),
                         card_w - Emu(240000), card_h - Emu(700000), lines)

    # =====================================================================
    # SLIDE 9 — Dashboard & Knowledge (screenshot placeholders)
    # =====================================================================
    slide = prs.slides.add_slide(content1)
    set_text(slide.placeholders[13], "INNOVATION 4",
             font_name=FONT_MED, font_size=Pt(10), color=ROYAL)
    set_text(slide.placeholders[0], "Dashboard + Knowledge Base",
             font_name=FONT_MED, font_size=Pt(36), color=NAVY)
    for idx in [2, 14]:
        if idx in slide.placeholders:
            clear_placeholder(slide.placeholders[idx])

    # Two screenshot placeholders side by side
    ph_w = Emu(5600000)
    ph_h = Emu(3800000)
    ph_y = Emu(2200000)
    gap = Emu(300000)
    total_ph = 2 * ph_w + gap
    phx = (SLIDE_W - total_ph) // 2

    add_screenshot_placeholder(slide, phx, ph_y, ph_w, ph_h, "Project Dashboard")
    add_screenshot_placeholder(slide, phx + ph_w + gap, ph_y, ph_w, ph_h, "Knowledge Base")

    # Caption below
    add_rich_textbox(slide, Emu(300000), Emu(6100000), Emu(11500000), Emu(500000), [
        {"text": "Pipeline view  \u2022  Doughnut charts  \u2022  Team burnup  \u2022  Action item tracking  \u2022  Searchable knowledge entries",
         "font_name": FONT, "font_size": Pt(11), "color": DARK_GRAY,
         "alignment": PP_ALIGN.CENTER},
    ])

    # =====================================================================
    # SLIDE 10 — Architecture (visual diagram)
    # =====================================================================
    slide = prs.slides.add_slide(content1)
    set_text(slide.placeholders[13], "ARCHITECTURE",
             font_name=FONT_MED, font_size=Pt(10), color=ROYAL)
    set_text(slide.placeholders[0], "How It's Built",
             font_name=FONT_MED, font_size=Pt(36), color=NAVY)
    for idx in [2, 14]:
        if idx in slide.placeholders:
            clear_placeholder(slide.placeholders[idx])

    # Architecture: 4 horizontal layers
    layer_x = Emu(300000)
    layer_w = Emu(11500000)
    layer_h = Emu(900000)
    layer_gap = Emu(120000)
    start_y = Emu(2000000)

    layers = [
        ("EXTERNAL SERVICES", "Jira Cloud  |  Confluence  |  Zoom  |  Claude / Gemini / Ollama", ROYAL, WHITE, SKY),
        ("API CONNECTORS", "Auth + Retry + Pagination  |  Rate limiting  |  Error handling", NAVY, WHITE, SKY),
        ("CORE ENGINE", "7 LLM Agents  |  Approval Engine  |  Orchestrator  |  ProjectContextService", AMBER, NAVY, NAVY),
        ("WEB FRONTEND", "FastAPI + HTMX + Jinja2  |  Chart.js  |  10 page views  |  Dark mode", GREEN, WHITE, WHITE),
    ]

    for i, (title, desc, fill, title_c, desc_c) in enumerate(layers):
        y = start_y + i * (layer_h + layer_gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, layer_x, y, layer_w, layer_h)
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.fill.background()
        box.adjustments[0] = 0.08

        add_rich_textbox(slide, layer_x + Emu(200000), y + Emu(80000),
                         layer_w - Emu(400000), layer_h - Emu(160000), [
            {"text": title, "font_name": FONT_MED, "font_size": Pt(14),
             "color": title_c, "bold": True},
            {"text": desc, "font_name": FONT, "font_size": Pt(12),
             "color": desc_c, "space_before": Pt(4)},
        ])

    # Down arrows between layers
    for i in range(3):
        y_top = start_y + (i + 1) * (layer_h + layer_gap) - layer_gap
        mid_x = SLIDE_W // 2
        add_arrow(slide, mid_x, y_top, mid_x, y_top + layer_gap, color=LIGHT_GRAY, width=Pt(2))

    # Side annotation: SQLite + .env
    add_rounded_box(slide, Emu(300000), start_y + 4 * (layer_h + layer_gap) + Emu(100000),
                    Emu(5500000), Emu(600000), OFF_WHITE,
                    border_color=LIGHT_GRAY,
                    text="SQLite (13 tables, WAL mode)  |  .env config  |  TTL cache",
                    font_name=FONT, font_size=Pt(11), text_color=DARK_GRAY)

    # Tech stack badge
    add_rounded_box(slide, Emu(6100000), start_y + 4 * (layer_h + layer_gap) + Emu(100000),
                    Emu(5700000), Emu(600000), OFF_WHITE,
                    border_color=LIGHT_GRAY,
                    text="Python 3.12+  |  200+ tests  |  80+ files  |  ~15,000 LOC",
                    font_name=FONT, font_size=Pt(11), text_color=DARK_GRAY)

    # =====================================================================
    # SLIDE 11 — Built with AI (meta story)
    # =====================================================================
    slide = prs.slides.add_slide(color_navy)
    set_text(slide.placeholders[13], "META",
             font_name=FONT_MED, font_size=Pt(10), color=SKY)
    set_text(slide.placeholders[0], "Built with AI,\nfor AI Governance",
             font_name=FONT_MED, font_size=Pt(36), color=WHITE)
    for idx in [2, 14]:
        if idx in slide.placeholders:
            clear_placeholder(slide.placeholders[idx])

    # Big stat boxes
    stat_w = Emu(3400000)
    stat_h = Emu(1400000)
    stat_y = Emu(2400000)
    stat_gap = Emu(300000)
    total_stat = 3 * stat_w + 2 * stat_gap
    stat_x = (SLIDE_W - total_stat) // 2

    stats = [
        ("~15K", "lines of Python", "across 80+ files"),
        ("200+", "automated tests", "connectors, agents, routes"),
        ("7", "AI agents", "provider-agnostic LLM layer"),
    ]
    for i, (big, sub1, sub2) in enumerate(stats):
        x = stat_x + i * (stat_w + stat_gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, stat_y, stat_w, stat_h)
        box.fill.solid()
        box.fill.fore_color.rgb = ROYAL
        box.line.fill.background()
        box.adjustments[0] = 0.08
        add_rich_textbox(slide, x + Emu(80000), stat_y + Emu(80000),
                         stat_w - Emu(160000), stat_h - Emu(160000), [
            {"text": big, "font_name": FONT_MED, "font_size": Pt(44),
             "color": WHITE, "bold": True, "alignment": PP_ALIGN.CENTER},
            {"text": sub1, "font_name": FONT_MED, "font_size": Pt(14),
             "color": SKY, "alignment": PP_ALIGN.CENTER, "space_before": Pt(4)},
            {"text": sub2, "font_name": FONT, "font_size": Pt(11),
             "color": SKY, "alignment": PP_ALIGN.CENTER, "space_before": Pt(2)},
        ])

    # Bottom text
    add_rich_textbox(slide, Emu(500000), Emu(4200000), Emu(11000000), Emu(2200000), [
        {"text": "Entire codebase developed with Claude Code",
         "font_name": FONT_MED, "font_size": Pt(18), "color": WHITE,
         "bold": True, "space_after": Pt(8)},
        {"text": "Project Seat is both an AI product and an AI development case study. "
         "It demonstrates that AI-assisted engineering can produce production-quality, "
         "well-tested, well-documented software at unprecedented speed.",
         "font_name": FONT, "font_size": Pt(13), "color": SKY, "space_before": Pt(8)},
        {"text": "Provider-agnostic: Claude Sonnet  \u2022  Gemini 2.5 Flash  \u2022  Ollama (local)",
         "font_name": FONT, "font_size": Pt(12), "color": MED_GRAY, "space_before": Pt(16)},
    ])

    # =====================================================================
    # SLIDE 12 — Value Proposition (big visual)
    # =====================================================================
    slide = prs.slides.add_slide(bg_royal)
    set_text(slide.placeholders[10], "VALUE",
             font_name=FONT_MED, font_size=Pt(10), color=SKY)

    # Hero stat
    add_textbox(slide, Emu(500000), Emu(800000), Emu(11000000), Emu(1200000),
                "Hours \u2192 Minutes",
                font_name=FONT_MED, font_size=Pt(56), color=WHITE, bold=True)
    add_textbox(slide, Emu(500000), Emu(1900000), Emu(11000000), Emu(500000),
                "Meeting analysis, status reports, and health reviews — done in minutes, not hours",
                font_name=FONT, font_size=Pt(16), color=SKY)

    # Three value pillars
    pill_w = Emu(3400000)
    pill_h = Emu(3000000)
    pill_y = Emu(3000000)
    pill_gap = Emu(300000)
    total_pill = 3 * pill_w + 2 * pill_gap
    pill_x = (SLIDE_W - total_pill) // 2

    pillars = [
        ("Zero\nKnowledge Loss",
         "Every meeting transcript produces\nstructured, traceable outputs.\n\nRisks, decisions, and action items\ngo directly into Jira — nothing\nfalls through the cracks."),
        ("Regulatory\nReady",
         "Human-in-the-loop approval queue\nwith full audit trail.\n\nEvery AI action is explicitly\napproved, meeting IEC 62304 and\nISO 14971 traceability."),
        ("Works\nToday",
         "Connects to your existing\nAtlassian and Zoom stack.\n\nNo new infrastructure. Runs locally.\nProvider-agnostic LLM — Claude,\nGemini, or on-premises Ollama."),
    ]
    for i, (title, desc) in enumerate(pillars):
        x = pill_x + i * (pill_w + pill_gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, pill_y, pill_w, pill_h)
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY
        box.line.fill.background()
        box.adjustments[0] = 0.06

        add_rich_textbox(slide, x + Emu(150000), pill_y + Emu(150000),
                         pill_w - Emu(300000), pill_h - Emu(300000), [
            {"text": title, "font_name": FONT_MED, "font_size": Pt(22),
             "color": WHITE, "bold": True, "space_after": Pt(12)},
            {"text": desc, "font_name": FONT, "font_size": Pt(12),
             "color": SKY, "space_before": Pt(8)},
        ])

    # =====================================================================
    # SLIDE 13 — Thank You
    # =====================================================================
    slide = prs.slides.add_slide(cover1)
    for ph in slide.placeholders:
        clear_placeholder(ph)
    add_textbox(slide, Emu(505141), Emu(2400000), Emu(8500000), Emu(1500000),
                "THANK YOU",
                font_name=FONT_MED, font_size=Pt(60), color=NAVY, bold=True)
    add_rich_textbox(slide, Emu(505141), Emu(4100000), Emu(8500000), Emu(1500000), [
        {"text": "github.com/Renneberg1/Project-Seat",
         "font_name": FONT, "font_size": Pt(18), "color": ROYAL},
        {"text": "Questions?",
         "font_name": FONT_MED, "font_size": Pt(24), "color": NAVY,
         "space_before": Pt(24), "bold": True},
    ])
    set_text(slide.placeholders[13], "THOMAS RENNEBERG  |  MARCH 2026",
             font_name=FONT, font_size=Pt(10), color=MED_GRAY)

    # ── Save ──────────────────────────────────────────────────────────────
    output_path = os.path.join(os.path.dirname(__file__), "..", "Project Seat Presentation.pptx")
    prs.save(output_path)
    print(f"Saved: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
